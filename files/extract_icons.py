"""
extract_icons.py — достаёт векторные иконки из .pptx-шаблона (например
SberF1) в библиотеку для генератора + рендерит пронумерованный контактный
лист, по которому удобно выбирать нужные иконки.

Иконки в шаблоне — это сгруппированные freeform-фигуры (custom geometry).
Скрипт собирает все верхнеуровневые группы со слайдов-«наборов иконок».

Запуск:
    # 1) распаковать шаблон
    python -m pptx_generator.tools.extract_icons \
        --template SberF1_Шаблон.pptx \
        --slides 6 7 8 9 10 11 12 \
        --out-json full_icons.json \
        --out-sheet icon_sheet.pptx

Потом открой icon_sheet.pptx (или сконвертируй в PDF), посмотри номера и
собери из full_icons.json нужные id в assets/icons_sberf1.json в формате:
    {"icons": {"i73": {...}}, "hints": {"warning": "i73"}}
"""
from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path

from lxml import etree

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def q(tag: str) -> str:
    pfx, name = tag.split(":")
    return f"{{{A if pfx == 'a' else P}}}{name}"


def extract(template: str, slides: list[int]) -> list[dict]:
    """Собирает иконки (верхнеуровневые группы) с указанных слайдов."""
    icons = []
    with zipfile.ZipFile(template) as z:
        for sn in slides:
            data = z.read(f"ppt/slides/slide{sn}.xml")
            root = etree.fromstring(data)
            spTree = root.find(".//" + q("p:cSld") + "/" + q("p:spTree"))
            if spTree is None:
                continue
            for grp in spTree.findall(q("p:grpSp")):
                xfrm = grp.find(q("p:grpSpPr") + "/" + q("a:xfrm"))
                if xfrm is None:
                    continue
                chExt = xfrm.find(q("a:chExt"))
                cx = int(chExt.get("cx")) if chExt is not None else 1
                cy = int(chExt.get("cy")) if chExt is not None else 1
                icons.append({
                    "xml": etree.tostring(grp).decode(),
                    "aspect": (cy / cx) if cx else 1.0,
                })
    return icons


def render_sheet(icons: list[dict], out_pptx: str,
                 cols: int = 10, rows: int = 6) -> None:
    """Рендерит контактный лист с номерами для выбора иконок."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    per = cols * rows
    ICON = Emu(560000)
    x0, y0, dx, dy = 500000, 450000, 1150000, 1050000

    def recolor(g, hexc):
        for sf in g.iter(q("a:solidFill")):
            for ch in list(sf):
                sf.remove(ch)
            etree.SubElement(sf, q("a:srgbClr")).set("val", hexc)

    def stamp(slide, it, left, top, size, hexc):
        g = etree.fromstring(it["xml"])
        xfrm = g.find(q("p:grpSpPr") + "/" + q("a:xfrm"))
        xfrm.find(q("a:off")).set("x", str(left))
        xfrm.find(q("a:off")).set("y", str(top))
        xfrm.find(q("a:ext")).set("cx", str(size))
        xfrm.find(q("a:ext")).set("cy", str(int(size * it["aspect"])))
        recolor(g, hexc)
        slide.shapes._spTree.append(g)

    for page in range((len(icons) + per - 1) // per):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        bg.line.fill.background()
        for k in range(per):
            idx = page * per + k
            if idx >= len(icons):
                break
            c, r = k % cols, k // cols
            left = x0 + c * dx
            top = y0 + r * dy
            stamp(s, icons[idx], left, top, int(ICON), "111827")
            tb = s.shapes.add_textbox(Emu(left), Emu(top + int(ICON)),
                                      ICON, Emu(220000))
            p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(idx)
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor.from_string("0669E0")
    prs.save(out_pptx)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("--slides", nargs="+", type=int, required=True)
    ap.add_argument("--out-json", default="full_icons.json")
    ap.add_argument("--out-sheet", default="icon_sheet.pptx")
    args = ap.parse_args()

    icons = extract(args.template, args.slides)
    print(f"извлечено иконок: {len(icons)}")
    Path(args.out_json).write_text(json.dumps(icons, ensure_ascii=False))
    render_sheet(icons, args.out_sheet)
    print(f"библиотека: {args.out_json}\nконтактный лист: {args.out_sheet}")


if __name__ == "__main__":
    main()
