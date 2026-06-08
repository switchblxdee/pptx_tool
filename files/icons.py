"""
Иконки для дайджеста, нарисованные нативными средствами python-pptx.

Не используем картинки/эмодзи — рисуем простые векторные иконки через
шейпы и freeform. Это надёжно рендерится везде (PowerPoint, LibreOffice)
и масштабируется без потери качества.

Каждая функция принимает slide, позицию (left, top), размер (size) и цвет,
и добавляет иконку на слайд. Иконки квадратные (size × size).

Доступные иконки (по icon_hint из схемы):
- signal      — сигнал/волны
- lens        — лупа (поиск/анализ)
- arrow_up    — рост
- arrow_down  — снижение
- info        — информация
- warning     — внимание
- check       — готово
- clock       — время
- dot         — простая точка-маркер (fallback)
"""
from __future__ import annotations

from typing import Callable, Dict

from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _fill(shape, hex_color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)
    shape.line.fill.background()


def _outline(shape, hex_color: str, width_pt: float = 1.5) -> None:
    shape.fill.background()
    shape.line.color.rgb = _rgb(hex_color)
    shape.line.width = Pt(width_pt)


# --------------------------------------------------------------------------- #
# Отдельные иконки
# --------------------------------------------------------------------------- #

def icon_signal(slide, left, top, size, color: str) -> None:
    """Сигнал — три концентрические дуги (как Wi-Fi/звук)."""
    # Рисуем как точка + две дуги. Упрощённо — три круга разного размера.
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Emu(int(size * 0.38)), top + Emu(int(size * 0.6)),
        Emu(int(size * 0.24)), Emu(int(size * 0.24)),
    )
    _fill(dot, color)
    # Дуга (используем BLOCK_ARC недоступен везде — берём смайл-арку через CHORD)
    for i, scale in enumerate((0.55, 0.85)):
        arc = slide.shapes.add_shape(
            MSO_SHAPE.ARC,
            left + Emu(int(size * (0.5 - scale / 2))),
            top + Emu(int(size * (0.5 - scale / 2))),
            Emu(int(size * scale)), Emu(int(size * scale)),
        )
        _outline(arc, color, 2.0)


def icon_lens(slide, left, top, size, color: str) -> None:
    """Лупа — круг + ручка."""
    d = Emu(int(size * 0.6))
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Emu(int(size * 0.05)), top + Emu(int(size * 0.05)), d, d,
    )
    _outline(circle, color, 2.5)
    # Ручка — повёрнутый прямоугольник
    handle = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Emu(int(size * 0.55)), top + Emu(int(size * 0.55)),
        Emu(int(size * 0.35)), Emu(int(size * 0.12)),
    )
    _fill(handle, color)
    handle.rotation = 45


def icon_arrow_up(slide, left, top, size, color: str) -> None:
    """Стрелка вверх — рост."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.UP_ARROW,
        left + Emu(int(size * 0.2)), top + Emu(int(size * 0.1)),
        Emu(int(size * 0.6)), Emu(int(size * 0.8)),
    )
    _fill(arrow, color)


def icon_arrow_down(slide, left, top, size, color: str) -> None:
    """Стрелка вниз — снижение."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        left + Emu(int(size * 0.2)), top + Emu(int(size * 0.1)),
        Emu(int(size * 0.6)), Emu(int(size * 0.8)),
    )
    _fill(arrow, color)


def icon_info(slide, left, top, size, color: str) -> None:
    """Информация — круг с буквой i (точка + палочка)."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Emu(int(size * 0.9)), Emu(int(size * 0.9)),
    )
    _outline(circle, color, 2.0)
    # точка
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Emu(int(size * 0.4)), top + Emu(int(size * 0.2)),
        Emu(int(size * 0.1)), Emu(int(size * 0.1)),
    )
    _fill(dot, color)
    # палочка
    stem = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Emu(int(size * 0.41)), top + Emu(int(size * 0.38)),
        Emu(int(size * 0.08)), Emu(int(size * 0.32)),
    )
    _fill(stem, color)


def icon_warning(slide, left, top, size, color: str) -> None:
    """Внимание — треугольник с восклицательным знаком."""
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOCELES_TRIANGLE,
        left, top, Emu(int(size * 0.9)), Emu(int(size * 0.9)),
    )
    _fill(tri, color)


def icon_check(slide, left, top, size, color: str) -> None:
    """Галочка в круге — готово."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Emu(int(size * 0.9)), Emu(int(size * 0.9)),
    )
    _fill(circle, color)


def icon_clock(slide, left, top, size, color: str) -> None:
    """Часы — круг (стрелки опускаем для простоты/чистоты)."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Emu(int(size * 0.9)), Emu(int(size * 0.9)),
    )
    _outline(circle, color, 2.0)


def icon_dot(slide, left, top, size, color: str) -> None:
    """Простая точка-маркер (fallback)."""
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Emu(int(size * 0.3)), top + Emu(int(size * 0.3)),
        Emu(int(size * 0.4)), Emu(int(size * 0.4)),
    )
    _fill(dot, color)


# --------------------------------------------------------------------------- #
# Реестр и резолвер
# --------------------------------------------------------------------------- #

_ICON_REGISTRY: Dict[str, Callable] = {
    "signal": icon_signal,
    "lens": icon_lens,
    "search": icon_lens,
    "arrow_up": icon_arrow_up,
    "up": icon_arrow_up,
    "growth": icon_arrow_up,
    "arrow_down": icon_arrow_down,
    "down": icon_arrow_down,
    "info": icon_info,
    "warning": icon_warning,
    "alert": icon_warning,
    "check": icon_check,
    "done": icon_check,
    "clock": icon_clock,
    "time": icon_clock,
    "dot": icon_dot,
}


def draw_icon(slide, hint: str, left, top, size, color: str) -> None:
    """
    Рисует иконку по hint. Если hint неизвестен — рисует точку.

    :param hint: ключ иконки (signal, lens, arrow_up, ...)
    :param size: размер иконки в EMU (квадрат size × size)
    """
    if not hint:
        return
    fn = _ICON_REGISTRY.get(hint.lower().strip(), icon_dot)
    try:
        fn(slide, left, top, size, color)
    except Exception:
        # Иконки — украшение, не критичны. При сбое не валим рендер.
        pass


def has_icon(hint: str) -> bool:
    """Известна ли иконка с таким hint."""
    return bool(hint) and hint.lower().strip() in _ICON_REGISTRY
