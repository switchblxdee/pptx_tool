"""
Сборка дайджеста-презентации.

Реализует визуальный язык референса «Голос IT»:
- Cover-слайд: градиентный фон, крупный заголовок, pill-метки источников,
  4 белые KPI-карточки с лёгкой тенью, шапка с метаданными.
- Topic-слайды: градиентный фон, плашка с названием темы, 4 KPI-карточки
  в линию (с обводкой), список тем-проблем на пастельной карточке,
  бейджи `new`, pill-метки источников, футер с метаданными.

Технические решения:
- Градиент фона делается через прямой XML-патч background-fill
  (python-pptx не предоставляет высокоуровневого API для градиентов).
- Pill-метки — rounded rectangle с динамической шириной от длины текста.
- Все размеры — в Emu/Inches, единая константа LAYOUT-зоны.
"""
from __future__ import annotations

from pathlib import Path
from typing import List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .icons import draw_icon, has_icon
from .schemas import (
    AttentionSlide,
    ClosingSlide,
    CoverSlide,
    DigestSpec,
    DigestStyle,
    ExecutiveSummarySlide,
    KPICard,
    PatternsSlide,
    TopicItem,
    TopicSlide,
)

# --------------------------------------------------------------------------- #
# Константы
# --------------------------------------------------------------------------- #

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Зона безопасных отступов
MARGIN_X = Inches(0.55)

# Зона шапки/футера с метаданными
META_HEADER_TOP = Inches(0.25)
META_HEADER_HEIGHT = Inches(0.4)
FOOTER_BOTTOM = Inches(0.2)
FOOTER_HEIGHT = Inches(0.35)

# XML namespaces для прямого патчинга OOXML
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


# --------------------------------------------------------------------------- #
# Builder
# --------------------------------------------------------------------------- #

class DigestBuilder:
    """Собирает .pptx-дайджест по DigestSpec."""

    def __init__(self, spec: DigestSpec):
        self.spec = spec
        self.style = spec.style
        self.palette = spec.style.palette
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def build(self, output_path: str | Path) -> Path:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Считаем общее число слайдов для нумерации в футере
        total = 1  # cover
        if self.spec.executive_summary:
            total += 1
        total += len(self.spec.topics)
        if self.spec.patterns:
            total += 1
        if self.spec.attention:
            total += 1
        if self.spec.closing:
            total += 1

        page = 1

        # 1. Обложка
        cover_slide = self._new_slide()
        self._render_cover(cover_slide, self.spec.cover)
        self._render_meta_header(cover_slide)

        # 2. Executive summary (сразу после обложки)
        if self.spec.executive_summary:
            page += 1
            slide = self._new_slide()
            self._render_executive_summary(slide, self.spec.executive_summary)
            self._render_meta_footer(slide, page_number=page, total=total)

        # 3. Topic-слайды
        for topic in self.spec.topics:
            page += 1
            slide = self._new_slide()
            self._render_topic(slide, topic)
            self._render_meta_footer(slide, page_number=page, total=total)

        # 4. Сквозные паттерны
        if self.spec.patterns:
            page += 1
            slide = self._new_slide()
            self._render_patterns(slide, self.spec.patterns)
            self._render_meta_footer(slide, page_number=page, total=total)

        # 5. На что обратить внимание
        if self.spec.attention:
            page += 1
            slide = self._new_slide()
            self._render_attention(slide, self.spec.attention)
            self._render_meta_footer(slide, page_number=page, total=total)

        # 6. Итоги
        if self.spec.closing:
            page += 1
            slide = self._new_slide()
            self._render_closing(slide, self.spec.closing)
            self._render_meta_footer(slide, page_number=page, total=total)

        self.prs.save(output_path)
        return output_path

    def _new_slide(self):
        """Создаёт пустой слайд с градиентным фоном."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_gradient_background(slide)
        return slide

    # ----------------------------------------------------------------------- #
    # COVER
    # ----------------------------------------------------------------------- #

    def _render_cover(self, slide, cover: CoverSlide) -> None:
        # Заголовок («Голос IT»)
        self._add_text(
            slide, cover.title,
            left=Inches(0.8), top=Inches(2.0),
            width=Inches(11.5), height=Inches(0.9),
            font=self.style.typography.heading_font,
            size=44, bold=True, color=self.palette.text_dark,
        )

        # Подзаголовок («дайджест для руководства Блока T»)
        self._add_text(
            slide, cover.subtitle,
            left=Inches(0.8), top=Inches(2.85),
            width=Inches(11.5), height=Inches(0.9),
            font=self.style.typography.heading_font,
            size=32, bold=True, color=self.palette.text_dark,
        )

        # Описание (мелким текстом)
        if cover.description:
            self._add_text(
                slide, cover.description,
                left=Inches(0.8), top=Inches(3.85),
                width=Inches(11.5), height=Inches(0.7),
                font=self.style.typography.body_font,
                size=14, color=self.palette.text_muted,
            )

        # Pill-метки источников
        self._render_pills_row(
            slide, cover.source_tags,
            top=Inches(4.7),
            font_size=11,
            text_color=self.palette.text_dark,
            bg_color=self.palette.kpi_bg,
        )

        # KPI-карточки внизу
        self._render_kpi_row(
            slide, cover.kpis,
            top=Inches(5.5),
            height=Inches(1.6),
            value_size=36,
            label_size=11,
            value_color=self.palette.text_dark,
            label_color=self.palette.text_muted,
            card_bg=self.palette.kpi_bg,
            border_color=None,  # на cover карточки без обводки
        )

    # ----------------------------------------------------------------------- #
    # TOPIC
    # ----------------------------------------------------------------------- #

    def _render_topic(self, slide, topic: TopicSlide) -> None:
        # Зона заголовка темы слева (плашка как на референсе)
        title_box_left = Inches(0.6)
        title_box_top = Inches(0.95)
        title_box_width = Inches(4.2)
        title_box_height = Inches(1.5)

        # Адаптивный размер шрифта по длине названия темы
        if len(topic.title) > 30:
            title_size = 24
        elif len(topic.title) > 18:
            title_size = 28
        else:
            title_size = 32

        # Чтобы тема визуально «звучала» — лёгкий блок-плашка с акцентом
        self._add_text(
            slide, topic.title,
            left=title_box_left, top=title_box_top,
            width=title_box_width, height=title_box_height,
            font=self.style.typography.heading_font,
            size=title_size, bold=True, color=self.palette.text_dark,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # KPI-карточки в линию справа от заголовка темы
        kpi_left_start = Inches(5.0)
        kpi_total_width = Inches(7.8)
        kpi_gap = Inches(0.2)
        n_kpi = len(topic.kpis)
        kpi_card_width = Emu(int((kpi_total_width - kpi_gap * (n_kpi - 1)) / n_kpi))

        for i, kpi in enumerate(topic.kpis):
            left = kpi_left_start + (kpi_card_width + kpi_gap) * i
            self._render_kpi_card(
                slide, kpi,
                left=left, top=Inches(0.95),
                width=kpi_card_width, height=Inches(1.5),
                value_size=28,
                label_size=10,
                value_color=self.palette.text_dark,
                label_color=self.palette.text_muted,
                card_bg=None,                       # прозрачный (на градиенте)
                border_color=self.palette.text_dark,  # с обводкой
            )

        # Карточка-список тем (пастельный фон)
        list_top = Inches(2.7)
        list_height = Inches(3.5)
        self._render_topic_items_card(
            slide, topic.items,
            left=MARGIN_X, top=list_top,
            width=SLIDE_WIDTH - MARGIN_X * 2,
            height=list_height,
        )

        # Pill-метки источников внизу
        if topic.source_tags:
            self._add_text(
                slide, "Источники:",
                left=MARGIN_X, top=Inches(6.4),
                width=Inches(1.2), height=Inches(0.3),
                font=self.style.typography.body_font,
                size=11, color=self.palette.text_muted,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            self._render_pills_row(
                slide, topic.source_tags,
                top=Inches(6.4),
                start_left=Inches(1.65),
                font_size=10,
                text_color=self.palette.text_dark,
                bg_color=self.palette.kpi_bg,
            )

    def _render_topic_items_card(
        self, slide, items: List[TopicItem],
        left, top, width, height,
    ) -> None:
        """Карточка с пастельным фоном, содержащая список тем-проблем."""
        # Фон карточки
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.adjustments[0] = 0.04
        self._fill_solid(card, self.palette.card_bg)
        card.line.fill.background()
        # Лёгкая тень (без неё карточка плоская)
        self._apply_subtle_shadow(card)

        # Раскладка items внутри карточки
        n = len(items)
        inner_pad_x = Inches(0.5)
        inner_pad_y = Inches(0.3)

        # Высота каждой строки = (общая высота - отступы) / n
        usable_height = height - inner_pad_y * 2
        row_height = Emu(int(usable_height / n))

        # Колонки внутри строки:
        # [номер|заголовок+цитата .............. период | mentions]
        col_num_width = Inches(0.4)
        col_period_width = Inches(2.0)
        col_mentions_width = Inches(0.9)
        col_content_left = left + inner_pad_x + col_num_width + Inches(0.1)
        col_content_width = (
            width - inner_pad_x * 2 - col_num_width - col_period_width - col_mentions_width - Inches(0.3)
        )
        col_period_left = left + width - inner_pad_x - col_mentions_width - col_period_width - Inches(0.1)
        col_mentions_left = left + width - inner_pad_x - col_mentions_width

        # Максимум упоминаний — для нормировки мини-баров
        max_mentions = max((it.mentions for it in items), default=0) or 1

        for i, item in enumerate(items):
            row_top = top + inner_pad_y + row_height * i

            # Номер
            self._add_text(
                slide, f"{i + 1}.",
                left=left + inner_pad_x, top=row_top,
                width=col_num_width, height=Inches(0.4),
                font=self.style.typography.body_font,
                size=14, bold=True, color=self.palette.text_dark,
                anchor=MSO_ANCHOR.TOP,
            )

            # Бейдж `new` (если есть) — рисуем ПОСЛЕ заголовка,
            # чтобы он не наезжал. Заголовок сжимаем по ширине.
            badge_width = Inches(0.5) if item.is_new else Inches(0)
            content_left = col_content_left
            # При наличии бейджа отступаем заголовок вправо на ширину бейджа
            if item.is_new:
                self._render_new_badge(
                    slide,
                    left=col_content_left,
                    top=row_top + Inches(0.05),
                )
                content_left = col_content_left + badge_width + Inches(0.1)
            content_width = col_content_width - (
                badge_width + Inches(0.1) if item.is_new else Emu(0)
            )

            # Заголовок темы
            self._add_text(
                slide, item.title,
                left=content_left, top=row_top,
                width=content_width, height=Inches(0.4),
                font=self.style.typography.heading_font,
                size=15, bold=True, color=self.palette.text_dark,
            )

            # Цитата под заголовком (курсив) — начинается с того же left,
            # что и заголовок (если бейдж был — после него)
            if item.quote:
                self._add_text(
                    slide, f"«{item.quote}»",
                    left=col_content_left, top=row_top + Inches(0.4),
                    width=col_content_width, height=Inches(0.6),
                    font=self.style.typography.body_font,
                    size=11, italic=True, color=self.palette.text_muted,
                )

            # Период (справа)
            self._add_text(
                slide, item.period,
                left=col_period_left, top=row_top + Inches(0.05),
                width=col_period_width, height=Inches(0.4),
                font=self.style.typography.body_font,
                size=12, color=self.palette.text_dark,
                align=PP_ALIGN.RIGHT,
            )

            # Mentions (число + подпись)
            self._add_text(
                slide, str(item.mentions),
                left=col_mentions_left, top=row_top,
                width=col_mentions_width, height=Inches(0.4),
                font=self.style.typography.heading_font,
                size=20, bold=True, color=self.palette.text_dark,
                align=PP_ALIGN.RIGHT,
            )
            self._add_text(
                slide, "упом.",
                left=col_mentions_left, top=row_top + Inches(0.4),
                width=col_mentions_width, height=Inches(0.3),
                font=self.style.typography.body_font,
                size=10, color=self.palette.text_muted,
                align=PP_ALIGN.RIGHT,
            )

            # Мини-бар под числом упоминаний — визуализация относительного веса.
            # Длина пропорциональна mentions / max_mentions.
            bar_full_width = col_mentions_width
            bar_ratio = item.mentions / max_mentions
            bar_width = Emu(max(int(bar_full_width * bar_ratio), Inches(0.1)))
            bar_top = row_top + Inches(0.72)
            # Фоновая дорожка (бледная)
            track = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                col_mentions_left, bar_top, bar_full_width, Inches(0.07),
            )
            track.adjustments[0] = 0.5
            self._fill_solid(track, self.palette.text_muted)
            track.line.fill.background()
            track.fill.fore_color.rgb = self._rgb(self.palette.text_muted)
            # Заполнение (акцентный цвет), выровнено по правому краю
            fill_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                col_mentions_left + bar_full_width - bar_width, bar_top,
                bar_width, Inches(0.07),
            )
            fill_bar.adjustments[0] = 0.5
            self._fill_solid(fill_bar, self.palette.accent)
            fill_bar.line.fill.background()

            # Разделитель между строками (кроме последней)
            if i < n - 1:
                sep_y = row_top + row_height - Inches(0.05)
                separator = slide.shapes.add_connector(
                    1,  # MSO_CONNECTOR.STRAIGHT
                    left + inner_pad_x, sep_y,
                    left + width - inner_pad_x, sep_y,
                )
                separator.line.color.rgb = self._rgb(self.palette.text_muted)
                separator.line.width = Pt(0.5)

    def _render_new_badge(self, slide, left, top) -> None:
        """Оранжевый бейдж `new` как на референсе."""
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top,
            Inches(0.4), Inches(0.22),
        )
        badge.adjustments[0] = 0.5
        self._fill_solid(badge, self.palette.badge)
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = "new"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = self.style.typography.body_font
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = self._rgb("FFFFFF")

    # ----------------------------------------------------------------------- #
    # АНАЛИТИЧЕСКИЕ СЛАЙДЫ
    # ----------------------------------------------------------------------- #

    def _render_slide_title(self, slide, title: str, accent_bar: bool = True) -> None:
        """Общий заголовок аналитического слайда с акцентным маркером."""
        if accent_bar:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.6), Inches(0.7),
                Inches(0.1), Inches(0.6),
            )
            self._fill_solid(bar, self.palette.accent)
            bar.line.fill.background()

        self._add_text(
            slide, title,
            left=Inches(0.85), top=Inches(0.6),
            width=Inches(11.5), height=Inches(0.8),
            font=self.style.typography.heading_font,
            size=32, bold=True, color=self.palette.text_dark,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def _render_executive_summary(self, slide, s: ExecutiveSummarySlide) -> None:
        """Executive summary: вводный абзац + крупные тезисы."""
        self._render_slide_title(slide, s.title)

        top = Inches(1.7)

        # Вводный абзац на пастельной плашке
        if s.intro:
            intro_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                MARGIN_X, top, SLIDE_WIDTH - MARGIN_X * 2, Inches(0.9),
            )
            intro_card.adjustments[0] = 0.06
            self._fill_solid(intro_card, self.palette.card_bg)
            intro_card.line.fill.background()
            self._apply_subtle_shadow(intro_card)

            self._add_text(
                slide, s.intro,
                left=MARGIN_X + Inches(0.35), top=top + Inches(0.1),
                width=SLIDE_WIDTH - MARGIN_X * 2 - Inches(0.7), height=Inches(0.7),
                font=self.style.typography.body_font,
                size=14, color=self.palette.text_dark,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            top = top + Inches(1.2)

        # Тезисы: каждый — headline крупно + detail мельче
        n = len(s.points)
        usable_height = SLIDE_HEIGHT - top - Inches(0.7)
        row_height = Emu(int(usable_height / n))

        for i, point in enumerate(s.points):
            row_top = top + row_height * i

            # Номерной акцент-кружок
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                MARGIN_X, row_top + Inches(0.05),
                Inches(0.45), Inches(0.45),
            )
            self._fill_solid(num_circle, self.palette.accent)
            num_circle.line.fill.background()
            tf = num_circle.text_frame
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = str(i + 1)
            pp = tf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            for run in pp.runs:
                run.font.name = self.style.typography.heading_font
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = self._rgb(self.palette.kpi_bg)

            content_left = MARGIN_X + Inches(0.7)
            content_width = SLIDE_WIDTH - content_left - MARGIN_X

            # Headline
            self._add_text(
                slide, point.headline,
                left=content_left, top=row_top,
                width=content_width, height=Inches(0.45),
                font=self.style.typography.heading_font,
                size=18, bold=True, color=self.palette.text_dark,
            )
            # Detail
            if point.detail:
                self._add_text(
                    slide, point.detail,
                    left=content_left, top=row_top + Inches(0.45),
                    width=content_width, height=Inches(0.5),
                    font=self.style.typography.body_font,
                    size=13, color=self.palette.text_muted,
                )

    def _render_patterns(self, slide, s: PatternsSlide) -> None:
        """Сквозные паттерны: карточки-блоки с описанием закономерностей."""
        self._render_slide_title(slide, s.title)

        top = Inches(1.7)
        if s.intro:
            self._add_text(
                slide, s.intro,
                left=MARGIN_X, top=top,
                width=SLIDE_WIDTH - MARGIN_X * 2, height=Inches(0.6),
                font=self.style.typography.body_font,
                size=14, color=self.palette.text_muted,
            )
            top = top + Inches(0.75)

        n = len(s.patterns)
        usable_height = SLIDE_HEIGHT - top - Inches(0.7)
        gap = Inches(0.2)
        card_height = Emu(int((usable_height - gap * (n - 1)) / n))

        for i, pattern in enumerate(s.patterns):
            card_top = top + (card_height + gap) * i
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                MARGIN_X, card_top,
                SLIDE_WIDTH - MARGIN_X * 2, card_height,
            )
            card.adjustments[0] = 0.08
            self._fill_solid(card, self.palette.card_bg)
            card.line.fill.background()
            self._apply_subtle_shadow(card)

            inner_x = MARGIN_X + Inches(0.4)

            # Заголовок паттерна
            title_width = SLIDE_WIDTH - MARGIN_X * 2 - Inches(0.8) - Inches(1.6)
            self._add_text(
                slide, pattern.title,
                left=inner_x, top=card_top + Inches(0.15),
                width=title_width, height=Inches(0.4),
                font=self.style.typography.heading_font,
                size=16, bold=True, color=self.palette.text_dark,
            )
            # Описание
            self._add_text(
                slide, pattern.description,
                left=inner_x, top=card_top + Inches(0.55),
                width=title_width, height=card_height - Inches(0.6),
                font=self.style.typography.body_font,
                size=12, color=self.palette.text_muted,
            )

            # Бейдж «затронуто N тем» справа
            if pattern.affected_count is not None:
                badge_left = SLIDE_WIDTH - MARGIN_X - Inches(1.7)
                self._add_text(
                    slide, str(pattern.affected_count),
                    left=badge_left, top=card_top + Inches(0.15),
                    width=Inches(1.4), height=Inches(0.5),
                    font=self.style.typography.heading_font,
                    size=28, bold=True, color=self.palette.accent,
                    align=PP_ALIGN.RIGHT,
                )
                self._add_text(
                    slide, "тем затронуто",
                    left=badge_left, top=card_top + Inches(0.65),
                    width=Inches(1.4), height=Inches(0.3),
                    font=self.style.typography.body_font,
                    size=10, color=self.palette.text_muted,
                    align=PP_ALIGN.RIGHT,
                )

    def _render_attention(self, slide, s: AttentionSlide) -> None:
        """На что обратить внимание: список с цветовыми маркерами важности."""
        self._render_slide_title(slide, s.title)

        # Цвета severity
        severity_colors = {
            "высокий": self.palette.badge,
            "средний": self.palette.accent,
            "низкий": self.palette.text_muted,
        }

        top = Inches(1.8)
        n = len(s.items)
        gap = Inches(0.2)
        usable_height = SLIDE_HEIGHT - top - Inches(0.7)
        card_height = Emu(int((usable_height - gap * (n - 1)) / n))

        for i, item in enumerate(s.items):
            card_top = top + (card_height + gap) * i
            color = severity_colors.get(item.severity, self.palette.accent)

            # Карточка
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                MARGIN_X, card_top,
                SLIDE_WIDTH - MARGIN_X * 2, card_height,
            )
            card.adjustments[0] = 0.08
            self._fill_solid(card, self.palette.kpi_bg)
            card.line.fill.background()
            self._apply_subtle_shadow(card)

            # Цветная полоса severity слева
            sev_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                MARGIN_X + Inches(0.15), card_top + Inches(0.15),
                Inches(0.12), card_height - Inches(0.3),
            )
            sev_bar.adjustments[0] = 0.5
            self._fill_solid(sev_bar, color)
            sev_bar.line.fill.background()

            inner_x = MARGIN_X + Inches(0.5)
            content_width = SLIDE_WIDTH - inner_x - MARGIN_X - Inches(1.5)

            # Заголовок
            self._add_text(
                slide, item.title,
                left=inner_x, top=card_top + Inches(0.15),
                width=content_width, height=Inches(0.4),
                font=self.style.typography.heading_font,
                size=15, bold=True, color=self.palette.text_dark,
            )
            # Обоснование
            self._add_text(
                slide, item.rationale,
                left=inner_x, top=card_top + Inches(0.55),
                width=content_width, height=card_height - Inches(0.6),
                font=self.style.typography.body_font,
                size=12, color=self.palette.text_muted,
            )

            # Бейдж severity справа
            sev_badge_left = SLIDE_WIDTH - MARGIN_X - Inches(1.3)
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                sev_badge_left, card_top + (card_height - Inches(0.35)) / 2,
                Inches(1.1), Inches(0.35),
            )
            badge.adjustments[0] = 0.5
            self._fill_solid(badge, color)
            badge.line.fill.background()
            tf = badge.text_frame
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = item.severity
            pp = tf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            for run in pp.runs:
                run.font.name = self.style.typography.body_font
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = self._rgb(self.palette.kpi_bg)

    def _render_closing(self, slide, s: ClosingSlide) -> None:
        """Финальный слайд: обобщающий текст + опциональные итоговые KPI."""
        self._render_slide_title(slide, s.title)

        # Обобщающий текст на крупной плашке
        summary_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN_X, Inches(1.8),
            SLIDE_WIDTH - MARGIN_X * 2, Inches(2.2),
        )
        summary_card.adjustments[0] = 0.05
        self._fill_solid(summary_card, self.palette.card_bg)
        summary_card.line.fill.background()
        self._apply_subtle_shadow(summary_card)

        self._add_text(
            slide, s.summary,
            left=MARGIN_X + Inches(0.5), top=Inches(2.0),
            width=SLIDE_WIDTH - MARGIN_X * 2 - Inches(1.0), height=Inches(1.8),
            font=self.style.typography.body_font,
            size=16, color=self.palette.text_dark,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Итоговые KPI, если есть
        if s.kpis:
            self._render_kpi_row(
                slide, s.kpis,
                top=Inches(4.6),
                height=Inches(1.6),
                value_size=40, label_size=12,
                value_color=self.palette.text_dark,
                label_color=self.palette.text_muted,
                card_bg=self.palette.kpi_bg,
                border_color=None,
            )

    # ----------------------------------------------------------------------- #
    # KPI карточки (одинаковая логика для cover и topic)
    # ----------------------------------------------------------------------- #

    def _render_kpi_row(
        self, slide, kpis: List[KPICard],
        top, height,
        value_size: int, label_size: int,
        value_color: str, label_color: str,
        card_bg: Optional[str] = None,
        border_color: Optional[str] = None,
    ) -> None:
        """Ряд KPI-карточек по всей ширине слайда."""
        n = len(kpis)
        total_width = SLIDE_WIDTH - MARGIN_X * 2
        gap = Inches(0.25)
        card_width = Emu(int((total_width - gap * (n - 1)) / n))

        for i, kpi in enumerate(kpis):
            left = MARGIN_X + (card_width + gap) * i
            self._render_kpi_card(
                slide, kpi,
                left=left, top=top,
                width=card_width, height=height,
                value_size=value_size, label_size=label_size,
                value_color=value_color, label_color=label_color,
                card_bg=card_bg, border_color=border_color,
            )

    def _render_kpi_card(
        self, slide, kpi: KPICard,
        left, top, width, height,
        value_size: int, label_size: int,
        value_color: str, label_color: str,
        card_bg: Optional[str], border_color: Optional[str],
    ) -> None:
        """Одна KPI-карточка: значение крупно, подпись мельче."""
        # Фон карточки
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.adjustments[0] = 0.08
        if card_bg:
            self._fill_solid(card, card_bg)
            card.line.fill.background()
            self._apply_subtle_shadow(card)
        else:
            # Прозрачный фон, только обводка
            card.fill.background()

        if border_color:
            card.line.color.rgb = self._rgb(border_color)
            card.line.width = Pt(1.0)

        # Иконка в правом верхнем углу (если задан icon_hint)
        icon_hint = getattr(kpi, "icon_hint", None)
        if icon_hint and has_icon(icon_hint):
            icon_size = Inches(0.28)
            icon_color = value_color if card_bg else self.palette.accent
            draw_icon(
                slide, icon_hint,
                left=left + width - icon_size - Inches(0.18),
                top=top + Inches(0.16),
                size=icon_size,
                color=icon_color,
            )

        # Значение
        self._add_text(
            slide, kpi.value,
            left=left, top=top + Inches(0.15),
            width=width, height=Emu(int(height * 0.55)),
            font=self.style.typography.heading_font,
            size=value_size, bold=True, color=value_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Подпись
        label_top = top + Emu(int(height * 0.6))
        self._add_text(
            slide, kpi.label,
            left=left + Inches(0.1), top=label_top,
            width=width - Inches(0.2), height=Emu(int(height * 0.35)),
            font=self.style.typography.body_font,
            size=label_size, color=label_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.TOP,
        )

    # ----------------------------------------------------------------------- #
    # Pill-метки
    # ----------------------------------------------------------------------- #

    def _render_pills_row(
        self, slide, tags: List[str],
        top,
        font_size: int,
        text_color: str,
        bg_color: str,
        start_left=None,
    ) -> None:
        """Ряд pill-меток с динамической шириной."""
        if start_left is None:
            start_left = MARGIN_X

        # Эвристика: ширина пилки = длина текста * шрифт * коэффициент + отступы
        # px_per_char ≈ font_size * 0.55 для русского текста
        # 1pt = 12700 EMU
        pill_height = Inches(0.36)
        pill_gap = Inches(0.15)
        inner_pad_x_emu = Inches(0.3)

        current_left = start_left
        for tag in tags:
            # Эвристика ширины pill. Кириллица шире латиницы, поэтому
            # коэффициент щедрый. Размер pt в EMU: 1pt ≈ 12700 EMU.
            # Эмпирически 0.75 даёт корректные ширины для смешанного текста.
            char_width_emu = int(font_size * 0.75 * 12700)
            pill_width = Emu(len(tag) * char_width_emu) + inner_pad_x_emu * 2

            # Если выходит за правую границу — не рисуем больше.
            if current_left + pill_width > SLIDE_WIDTH - MARGIN_X:
                break

            pill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                current_left, top, pill_width, pill_height,
            )
            pill.adjustments[0] = 0.5  # полная скругленность
            self._fill_solid(pill, bg_color)
            pill.line.fill.background()
            self._apply_subtle_shadow(pill)

            tf = pill.text_frame
            tf.margin_left = inner_pad_x_emu
            tf.margin_right = inner_pad_x_emu
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = tag
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = self.style.typography.body_font
                run.font.size = Pt(font_size)
                run.font.color.rgb = self._rgb(text_color)

            current_left = current_left + pill_width + pill_gap

    # ----------------------------------------------------------------------- #
    # Метаданные (шапка и футер)
    # ----------------------------------------------------------------------- #

    def _render_meta_header(self, slide) -> None:
        """Шапка с датой/периодом/номером выпуска (для cover-слайда)."""
        meta = self.spec.meta
        text = f"{meta.issue_date}    •    Период: {meta.period}    •    {meta.issue_number}"

        # Полупрозрачная плашка справа
        pill_width = Inches(7.5)
        pill_left = SLIDE_WIDTH - MARGIN_X - pill_width

        self._add_text(
            slide, text,
            left=pill_left, top=META_HEADER_TOP,
            width=pill_width, height=META_HEADER_HEIGHT,
            font=self.style.typography.body_font,
            size=11, color=self.palette.text_dark,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    def _render_meta_footer(self, slide, page_number: int, total: int) -> None:
        """Футер с расширенной мета-информацией (для topic-слайдов)."""
        meta = self.spec.meta
        parts = [meta.issue_date, f"Период: {meta.period}", meta.issue_number]
        if meta.next_issue:
            parts.append(f"Следующий выпуск: {meta.next_issue}")
        if meta.note:
            parts.append(meta.note)
        text = "    •    ".join(parts)

        # Слева — мета
        self._add_text(
            slide, text,
            left=MARGIN_X, top=SLIDE_HEIGHT - FOOTER_BOTTOM - FOOTER_HEIGHT,
            width=SLIDE_WIDTH - MARGIN_X * 2 - Inches(0.7),
            height=FOOTER_HEIGHT,
            font=self.style.typography.body_font,
            size=9, color=self.palette.text_muted,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Справа — номер страницы
        self._add_text(
            slide, f"стр. {page_number}",
            left=SLIDE_WIDTH - MARGIN_X - Inches(0.7),
            top=SLIDE_HEIGHT - FOOTER_BOTTOM - FOOTER_HEIGHT,
            width=Inches(0.7), height=FOOTER_HEIGHT,
            font=self.style.typography.body_font,
            size=9, color=self.palette.text_muted,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ----------------------------------------------------------------------- #
    # Градиентный фон (прямой XML)
    # ----------------------------------------------------------------------- #

    def _set_gradient_background(self, slide) -> None:
        """
        Устанавливает горизонтальный градиент на фон слайда.

        python-pptx не имеет API для градиента фона, поэтому патчим XML
        напрямую. Это стабильно работает и в PowerPoint, и в LibreOffice.
        """
        start_hex = self.palette.gradient_start.lstrip("#")
        end_hex = self.palette.gradient_end.lstrip("#")

        # Берём bg-элемент слайда и заменяем fill
        bg = slide.background
        # Создаём <p:bgPr> с градиентной заливкой
        gradient_xml = f"""
        <p:bgPr xmlns:p="{NS_P}" xmlns:a="{NS_A}">
            <a:gradFill flip="none" rotWithShape="1">
                <a:gsLst>
                    <a:gs pos="0">
                        <a:srgbClr val="{start_hex}"/>
                    </a:gs>
                    <a:gs pos="100000">
                        <a:srgbClr val="{end_hex}"/>
                    </a:gs>
                </a:gsLst>
                <a:lin ang="0" scaled="1"/>
            </a:gradFill>
        </p:bgPr>
        """

        # Заменяем существующий bg-element
        bg_elem = bg._element  # CT_Background
        new_bg_pr = etree.fromstring(gradient_xml)

        # Удаляем старый bgPr/bgRef если есть
        for tag in ("p:bgPr", "p:bgRef"):
            for old in bg_elem.findall(qn(tag)):
                bg_elem.remove(old)

        bg_elem.append(new_bg_pr)

    # ----------------------------------------------------------------------- #
    # Утилиты текста и стиля
    # ----------------------------------------------------------------------- #

    def _add_text(
        self, slide, text: str,
        left, top, width, height,
        font: str, size: int, color: str,
        bold: bool = False, italic: bool = False,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    ) -> None:
        """Единая точка создания текстового бокса."""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.text = text
        p = tf.paragraphs[0]
        p.alignment = align
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = self._rgb(color)

    def _apply_subtle_shadow(self, shape) -> None:
        """Лёгкая тень — выделяет карточки на градиентном фоне."""
        # python-pptx не имеет высокоуровневого API для shadow,
        # патчим напрямую через spPr/effectLst
        sp_pr = shape.fill._xPr  # spPr
        # Если уже есть effectLst — не дублируем
        if sp_pr.find(qn("a:effectLst")) is not None:
            return
        effect_xml = """
        <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:outerShdw blurRad="50800" dist="25400" dir="5400000" algn="t" rotWithShape="0">
                <a:srgbClr val="000000"><a:alpha val="15000"/></a:srgbClr>
            </a:outerShdw>
        </a:effectLst>
        """
        sp_pr.append(etree.fromstring(effect_xml))

    def _fill_solid(self, shape, hex_color: str) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(hex_color)

    @staticmethod
    def _rgb(hex_color: str) -> RGBColor:
        return RGBColor.from_string(hex_color.lstrip("#"))
