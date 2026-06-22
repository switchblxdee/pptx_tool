"""
brand_icons.py — иконки из фирменного шаблона SberF1.

Иконки в шаблоне нарисованы векторной геометрией (custGeom). Мы извлекли
полезный набор и сохранили как прозрачные PNG в двух вариантах — тёмном и
светлом — чтобы они читались и на светлых подложках, и на тёмных карточках.

Почему PNG, а не вектор: цель — надёжная отрисовка в Р7-Офис (OnlyOffice),
PowerPoint и LibreOffice одинаково. Картинки (`add_picture`) — самый
универсально поддерживаемый элемент и не зависят от движка автофигур.

Публичный API (совместим с вызовами builder.py):
    has_brand_icon(hint)  -> bool
    draw_brand_icon(slide, hint, left, top, size, color) -> bool
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

_ICON_DIR = Path(__file__).resolve().parent / "assets" / "icons"

# Смысловой хинт -> базовое имя извлечённой иконки.
# Покрывает и авто-хинты билдера (signal/team/tools/...), и прямые имена.
_ALIASES = {
    # авто-хинты из _AUTO_ICON
    "signal": "chart", "warning": "warning", "bell": "new", "growth": "growth",
    "team": "users", "tools": "integration", "clock": "calendar", "bug": "error",
    "search": "info", "money": "data", "email": "document", "security": "security",
    # прямые / синонимы
    "document": "document", "doc": "document", "report": "document",
    "tag": "tag", "label": "tag", "category": "tag",
    "rocket": "rocket", "launch": "rocket", "speed": "rocket",
    "chart": "chart", "pie": "chart", "share": "chart", "distribution": "chart",
    "integration": "integration", "network": "integration", "node": "integration",
    "data": "data", "calc": "data", "metrics": "data",
    "users": "users", "people": "users", "audience": "users", "user": "users",
    "calendar": "calendar", "schedule": "calendar", "period": "calendar",
    "time": "calendar", "date": "calendar",
    "check": "check", "done": "check", "ok": "check", "success": "check",
    "wifi": "wifi", "signal_wave": "wifi", "connect": "wifi",
    "error": "error", "x": "error", "fail": "error", "close": "error",
    "new": "new", "add": "new", "plus": "new",
    "info": "info", "information": "info",
    "alert": "warning", "attention": "warning", "risk": "warning",
    "list": "list", "topics": "list", "menu": "list", "items": "list",
    "up": "growth", "trend": "growth", "increase": "growth",
}


def _base(hint: Optional[str]) -> Optional[str]:
    if not hint:
        return None
    return _ALIASES.get(hint.lower().strip())


def _exists(base: str) -> bool:
    return (_ICON_DIR / f"icon_{base}_dark.png").exists()


def has_brand_icon(hint: Optional[str]) -> bool:
    b = _base(hint)
    return bool(b) and _exists(b)


def available_brand_hints() -> list[str]:
    return sorted(h for h in _ALIASES if _exists(_ALIASES[h]))


def _is_light(color: Optional[str]) -> bool:
    """Светлый ли запрошенный цвет иконки (по воспринимаемой яркости)."""
    if not color:
        return False
    s = str(color).lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        return False
    r, g, b = int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0 > 0.5


def draw_brand_icon(slide, hint, left, top, size, color) -> bool:
    """Кладёт брендовую иконку по хинту. Вариант (тёмный/светлый) выбирается
    по яркости запрошенного цвета: светлый цвет -> светлая иконка (для тёмных
    карточек) и наоборот. Возвращает True, если иконка вставлена."""
    base = _base(hint)
    if not base:
        return False
    variant = "light" if _is_light(color) else "dark"
    path = _ICON_DIR / f"icon_{base}_{variant}.png"
    if not path.exists():
        alt = _ICON_DIR / f"icon_{base}_{'dark' if variant == 'light' else 'light'}.png"
        path = alt if alt.exists() else path
    if not path.exists():
        return False
    try:
        from pptx.util import Emu
        slide.shapes.add_picture(str(path), Emu(int(left)), Emu(int(top)),
                                 Emu(int(size)), Emu(int(size)))
        return True
    except Exception:
        return False
