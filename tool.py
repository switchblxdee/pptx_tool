"""
LangChain tool для генерации PowerPoint-презентаций c полной кастомизацией.

Зависимости:
    pip install langchain-core pydantic python-pptx

Три уровня кастомизации (всё опционально, дефолты разумные):

  1) Глобально — PresentationStyle: палитра, размер слайда, дефолтные шрифты/размеры.
  2) Слайд — SlideStyle: фон конкретного слайда, стили его блоков, акцентный круг.
  3) Элемент — любое строковое поле слайда можно заменить на StyledText:
       "title": "Привет"
     либо
       "title": {"text": "Привет", "style": {"size": 60, "color": "FF0000"}}

  Финальный стиль = дефолты ← palette ← presentation.defaults ← slide.style ← element.style

Минимальный пример (без кастомизации):
    create_presentation.invoke({
        "title": "Демо",
        "slides": [
            {"type": "title", "title": "Привет", "subtitle": "мир"},
            {"type": "content", "title": "Пункты", "bullets": ["A", "B", "C"]},
        ],
        "output_path": "out.pptx",
    })

Полный пример с кастомизацией:
    create_presentation.invoke({
        "title": "Кастом",
        "style": {
            "slide_width_inches": 13.333,
            "slide_height_inches": 7.5,
            "custom_palette": {"primary": "FF6B35", "secondary": "F7C59F",
                               "accent": "FFFFFF", "bg": "FFFCF2",
                               "text": "1A1A1A", "muted": "888888"},
            "defaults": {
                "title_style": {"font": "Helvetica", "size": 40, "bold": True},
                "body_style":  {"font": "Helvetica", "size": 18},
            },
        },
        "slides": [
            {
                "type": "content",
                "title": {"text": "ВАЖНО", "style": {"size": 60, "color": "990011"}},
                "bullets": [
                    "Обычный пункт",
                    {"text": "Особый пункт", "style": {"italic": True, "bold": True}},
                ],
                "style": {
                    "bg_color": "1E2761",
                    "accent_dot": {"enabled": False},
                    "body_style": {"color": "FFFFFF"},
                },
            },
        ],
        "output_path": "custom.pptx",
    })
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Union

from langchain_core.tools import tool
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pydantic import BaseModel, Field


# ---------- Палитры ----------

PALETTES: Dict[str, Dict[str, str]] = {
    "midnight_executive": {
        "primary": "1E2761", "secondary": "CADCFC", "accent": "FFFFFF",
        "bg": "FFFFFF", "text": "1E2761", "muted": "6B7280",
    },
    "forest_moss": {
        "primary": "2C5F2D", "secondary": "97BC62", "accent": "F5F5F5",
        "bg": "FFFFFF", "text": "1F2937", "muted": "6B7280",
    },
    "coral_energy": {
        "primary": "F96167", "secondary": "F9E795", "accent": "2F3C7E",
        "bg": "FFFFFF", "text": "2F3C7E", "muted": "6B7280",
    },
    "warm_terracotta": {
        "primary": "B85042", "secondary": "E7E8D1", "accent": "A7BEAE",
        "bg": "FFFFFF", "text": "3B2F2F", "muted": "8C7A7A",
    },
    "ocean_gradient": {
        "primary": "065A82", "secondary": "1C7293", "accent": "21295C",
        "bg": "FFFFFF", "text": "21295C", "muted": "6B7280",
    },
    "charcoal_minimal": {
        "primary": "36454F", "secondary": "F2F2F2", "accent": "212121",
        "bg": "FFFFFF", "text": "212121", "muted": "6B7280",
    },
    "teal_trust": {
        "primary": "028090", "secondary": "00A896", "accent": "02C39A",
        "bg": "FFFFFF", "text": "1F2937", "muted": "6B7280",
    },
    "berry_cream": {
        "primary": "6D2E46", "secondary": "A26769", "accent": "ECE2D0",
        "bg": "FFFFFF", "text": "3B1F2B", "muted": "8C7A7A",
    },
    "cherry_bold": {
        "primary": "990011", "secondary": "FCF6F5", "accent": "2F3C7E",
        "bg": "FFFFFF", "text": "2F3C7E", "muted": "6B7280",
    },
}


# ---------- Стилевые модели ----------

AlignT = Literal["left", "center", "right"]
_ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


class TextStyle(BaseModel):
    """Стиль текстового блока. Все поля опциональны — задавай только нужное."""
    font: Optional[str] = Field(None, description="Имя шрифта, например 'Helvetica', 'Georgia'")
    size: Optional[float] = Field(None, description="Размер шрифта в пунктах")
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    color: Optional[str] = Field(None, description="HEX без #, например 'FF6B35'")
    align: Optional[AlignT] = Field(None, description="left | center | right")
    # «Стиль маркера»: для буллет-списков
    bullet_char: Optional[str] = Field(None, description="Символ маркера, по умолчанию '•'")
    line_spacing: Optional[float] = Field(None, description="Множитель межстрочного интервала, например 1.2")
    space_after_pt: Optional[float] = Field(None, description="Отступ после абзаца в пунктах")


class StyledText(BaseModel):
    """Текст с inline-стилем. Используется вместо обычной строки в любом поле."""
    text: str
    style: Optional[TextStyle] = None


# Тип для полей слайда: либо строка, либо словарь с text/style
TextField = Union[str, StyledText, Dict[str, Any]]
ListField = Union[List[str], List[Union[str, StyledText, Dict[str, Any]]]]


class AccentDotStyle(BaseModel):
    """Маленький цветной круг — визуальный мотив. По умолчанию включён."""
    enabled: bool = True
    color: Optional[str] = Field(None, description="HEX без #; по умолчанию primary/accent палитры")
    diameter_inches: float = 0.25
    left_inches: float = 0.6
    top_inches: float = 0.7


class SlideStyle(BaseModel):
    """Стиль конкретного слайда. Перекрывает PresentationStyle.defaults."""
    bg_color: Optional[str] = Field(None, description="HEX фона слайда")
    accent_dot: Optional[AccentDotStyle] = None

    # Стили для отдельных «ролей» внутри слайда. Какие применяются — зависит от типа слайда.
    title_style: Optional[TextStyle] = None
    subtitle_style: Optional[TextStyle] = None
    body_style: Optional[TextStyle] = None          # буллеты, обычный текст
    column_header_style: Optional[TextStyle] = None  # заголовки колонок в two_column
    stat_style: Optional[TextStyle] = None           # большая цифра в stat
    caption_style: Optional[TextStyle] = None        # подпись под цифрой


class PresentationStyle(BaseModel):
    """Глобальный стиль презентации."""
    slide_width_inches: float = 13.333
    slide_height_inches: float = 7.5
    custom_palette: Optional[Dict[str, str]] = Field(
        None,
        description=(
            "Полностью своя палитра. Должна содержать ключи: "
            "primary, secondary, accent, bg, text, muted. "
            "Если задана — переопределяет именованную palette."
        ),
    )
    defaults: Optional[SlideStyle] = Field(
        None,
        description="Дефолтные стили для всех слайдов колоды.",
    )


# ---------- Схемы слайдов ----------

class TitleSlide(BaseModel):
    type: Literal["title"]
    title: TextField
    subtitle: Optional[TextField] = None
    style: Optional[SlideStyle] = None


class ContentSlide(BaseModel):
    type: Literal["content"]
    title: TextField
    bullets: ListField
    style: Optional[SlideStyle] = None


class TwoColumnSlide(BaseModel):
    type: Literal["two_column"]
    title: TextField
    left_header: TextField
    left_bullets: ListField
    right_header: TextField
    right_bullets: ListField
    style: Optional[SlideStyle] = None


class StatSlide(BaseModel):
    type: Literal["stat"]
    title: TextField
    stat: TextField
    caption: TextField
    style: Optional[SlideStyle] = None


class SectionSlide(BaseModel):
    type: Literal["section"]
    title: TextField
    style: Optional[SlideStyle] = None


# ---------- Дефолтные стили (фолбэк, если ничего не задано) ----------

# Эти значения применяются последними, если ни презентация, ни слайд, ни элемент
# не указали соответствующее поле.
HARDCODED_DEFAULTS: Dict[str, TextStyle] = {
    "title": TextStyle(font="Georgia", size=36, bold=True, align="left"),
    "subtitle": TextStyle(font="Calibri", size=22, italic=True, align="left"),
    "body": TextStyle(font="Calibri", size=18, align="left",
                       bullet_char="•", space_after_pt=8),
    "column_header": TextStyle(font="Georgia", size=22, bold=True, align="left"),
    "stat": TextStyle(font="Georgia", size=160, bold=True, align="center"),
    "caption": TextStyle(font="Calibri", size=20, italic=True, align="center"),
    # Спецдефолты для конкретных типов слайдов:
    "title_on_title_slide": TextStyle(font="Georgia", size=54, bold=True, align="left"),
    "title_on_section_slide": TextStyle(font="Georgia", size=48, bold=True, align="left"),
}


# ---------- Утилиты ----------

def _hex_to_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#").upper())


def _coerce_styled_text(value: Any) -> StyledText:
    """Приводит строку или dict к StyledText."""
    if isinstance(value, StyledText):
        return value
    if isinstance(value, str):
        return StyledText(text=value)
    if isinstance(value, dict):
        return StyledText(**value)
    raise TypeError(f"Не удаётся привести к StyledText: {value!r}")


def _coerce_styled_list(items: List[Any]) -> List[StyledText]:
    return [_coerce_styled_text(it) for it in items]


def _merge_styles(*styles: Optional[TextStyle]) -> TextStyle:
    """Сливает стили слева направо: поздние перекрывают ранние, если не None."""
    out: Dict[str, Any] = {}
    for s in styles:
        if s is None:
            continue
        data = s.model_dump(exclude_none=True) if hasattr(s, "model_dump") else {
            k: v for k, v in s.__dict__.items() if v is not None
        }
        out.update(data)
    return TextStyle(**out)


def _resolve_color(color: Optional[str], fallback: str) -> str:
    return color if color else fallback


def _is_light_bg(hex_color: str) -> bool:
    """Определяет, светлый ли фон, по перцептивной яркости (формула sRGB luminance)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    # Простая аппроксимация, без gamma-correction — достаточно для решения "светлый/тёмный"
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return luminance > 0.5


def _on_bg_color(bg: str, palette: Dict[str, str]) -> str:
    """Выбирает контрастный цвет текста для заданного фона.
    На светлом фоне — palette['text'], на тёмном — palette['accent'] (или белый)."""
    if _is_light_bg(bg):
        return palette["text"]
    # На тёмном фоне используем accent, если он контрастный, иначе белый
    accent = palette.get("accent", "FFFFFF")
    return accent if _is_light_bg(accent) else "FFFFFF"


def _muted_on_bg(bg: str, palette: Dict[str, str]) -> str:
    """Приглушённый цвет для подписей, контрастный фону."""
    if _is_light_bg(bg):
        return palette.get("muted", "6B7280")
    return palette.get("secondary", "CCCCCC")


# ---------- Применение стилей при рендеринге ----------

def _apply_paragraph_style(p, style: TextStyle, default_color: str) -> None:
    if style.align is not None:
        p.alignment = _ALIGN_MAP[style.align]
    if style.space_after_pt is not None:
        p.space_after = Pt(style.space_after_pt)
    if style.line_spacing is not None:
        p.line_spacing = style.line_spacing


def _apply_run_style(run, style: TextStyle, default_color: str) -> None:
    run.font.name = style.font or "Calibri"
    if style.size is not None:
        run.font.size = Pt(style.size)
    if style.bold is not None:
        run.font.bold = style.bold
    if style.italic is not None:
        run.font.italic = style.italic
    run.font.color.rgb = _hex_to_rgb(style.color or default_color)


def _add_styled_textbox(
    slide, left, top, width, height,
    text: StyledText,
    base_style: TextStyle,
    default_color: str,
) -> None:
    """Текстовый блок: base_style + inline-стиль текста."""
    effective = _merge_styles(base_style, text.style)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    _apply_paragraph_style(p, effective, default_color)
    run = p.add_run()
    run.text = text.text
    _apply_run_style(run, effective, default_color)


def _add_styled_bullets(
    slide, left, top, width, height,
    items: List[StyledText],
    base_style: TextStyle,
    default_color: str,
) -> None:
    """Список буллетов; каждый элемент может иметь свой inline-стиль."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    bullet_char = base_style.bullet_char or "•"
    for i, item in enumerate(items):
        effective = _merge_styles(base_style, item.style)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _apply_paragraph_style(p, effective, default_color)
        run = p.add_run()
        char = effective.bullet_char or bullet_char
        run.text = f"{char}  {item.text}" if char else item.text
        _apply_run_style(run, effective, default_color)


def _set_slide_bg(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _add_accent_dot(slide, dot: AccentDotStyle, fallback_color: str) -> None:
    if not dot.enabled:
        return
    d = Inches(dot.diameter_inches)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(dot.left_inches), Inches(dot.top_inches), d, d,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(dot.color or fallback_color)
    shape.line.fill.background()


# ---------- Резолвер стилей: дефолт → presentation → slide ----------

class StyleResolver:
    """Собирает финальный стиль для каждой роли с учётом всех уровней."""

    def __init__(
        self,
        palette: Dict[str, str],
        presentation_defaults: Optional[SlideStyle],
        slide_style: Optional[SlideStyle],
    ):
        self.palette = palette
        self.pres = presentation_defaults or SlideStyle()
        self.slide = slide_style or SlideStyle()

    def bg_color(self, fallback: str) -> str:
        return self.slide.bg_color or self.pres.bg_color or fallback

    def accent_dot(self, default_color: str) -> AccentDotStyle:
        # Slide перекрывает presentation; если ни там, ни там — дефолтный круг
        dot = self.slide.accent_dot or self.pres.accent_dot or AccentDotStyle()
        if dot.color is None:
            dot = dot.model_copy(update={"color": default_color})
        return dot

    def style_for(self, role: str, hardcoded_key: str) -> TextStyle:
        """role — имя поля в SlideStyle (e.g. 'title_style'); hardcoded_key — ключ в HARDCODED_DEFAULTS."""
        return _merge_styles(
            HARDCODED_DEFAULTS[hardcoded_key],
            getattr(self.pres, role, None),
            getattr(self.slide, role, None),
        )


# ---------- Рендереры по типам слайдов ----------

def _render_title(slide, spec: TitleSlide, palette: Dict[str, str],
                   resolver: StyleResolver) -> None:
    bg = resolver.bg_color(palette["primary"])
    _set_slide_bg(slide, bg)
    _add_accent_dot(slide, resolver.accent_dot(palette["accent"]),
                    palette["accent"])

    title_style = resolver.style_for("title_style", "title_on_title_slide")
    default_title_color = _on_bg_color(bg, palette)
    _add_styled_textbox(
        slide, Inches(0.6), Inches(2.8), Inches(11.8), Inches(1.6),
        _coerce_styled_text(spec.title), title_style, default_title_color,
    )

    if spec.subtitle:
        sub_style = resolver.style_for("subtitle_style", "subtitle")
        default_sub_color = _muted_on_bg(bg, palette)
        _add_styled_textbox(
            slide, Inches(0.6), Inches(4.6), Inches(11.8), Inches(0.8),
            _coerce_styled_text(spec.subtitle), sub_style, default_sub_color,
        )


def _render_section(slide, spec: SectionSlide, palette: Dict[str, str],
                     resolver: StyleResolver) -> None:
    bg = resolver.bg_color(palette["primary"])
    _set_slide_bg(slide, bg)
    _add_accent_dot(slide, resolver.accent_dot(palette["accent"]),
                    palette["accent"])
    title_style = resolver.style_for("title_style", "title_on_section_slide")
    default_color = _on_bg_color(bg, palette)
    _add_styled_textbox(
        slide, Inches(0.6), Inches(3.0), Inches(11.8), Inches(1.5),
        _coerce_styled_text(spec.title), title_style, default_color,
    )


def _render_content(slide, spec: ContentSlide, palette: Dict[str, str],
                     resolver: StyleResolver) -> None:
    bg = resolver.bg_color(palette["bg"])
    _set_slide_bg(slide, bg)
    _add_accent_dot(slide, resolver.accent_dot(palette["primary"]),
                    palette["primary"])

    title_style = resolver.style_for("title_style", "title")
    default_title_color = _on_bg_color(bg, palette)
    _add_styled_textbox(
        slide, Inches(1.0), Inches(0.55), Inches(11.5), Inches(0.8),
        _coerce_styled_text(spec.title), title_style, default_title_color,
    )

    body_style = resolver.style_for("body_style", "body")
    default_body_color = _on_bg_color(bg, palette)
    _add_styled_bullets(
        slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0),
        _coerce_styled_list(spec.bullets), body_style, default_body_color,
    )


def _render_two_column(slide, spec: TwoColumnSlide, palette: Dict[str, str],
                        resolver: StyleResolver) -> None:
    bg = resolver.bg_color(palette["bg"])
    _set_slide_bg(slide, bg)
    _add_accent_dot(slide, resolver.accent_dot(palette["primary"]),
                    palette["primary"])

    title_style = resolver.style_for("title_style", "title")
    text_color = _on_bg_color(bg, palette)
    _add_styled_textbox(
        slide, Inches(1.0), Inches(0.55), Inches(11.5), Inches(0.8),
        _coerce_styled_text(spec.title), title_style, text_color,
    )

    header_style = resolver.style_for("column_header_style", "column_header")
    body_style = resolver.style_for("body_style", "body")

    # Левая колонка
    _add_styled_textbox(
        slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.6),
        _coerce_styled_text(spec.left_header), header_style, palette["primary"],
    )
    _add_styled_bullets(
        slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.5),
        _coerce_styled_list(spec.left_bullets), body_style, text_color,
    )
    # Правая колонка
    _add_styled_textbox(
        slide, Inches(6.9), Inches(1.8), Inches(5.5), Inches(0.6),
        _coerce_styled_text(spec.right_header), header_style, palette["primary"],
    )
    _add_styled_bullets(
        slide, Inches(6.9), Inches(2.5), Inches(5.5), Inches(4.5),
        _coerce_styled_list(spec.right_bullets), body_style, text_color,
    )


def _render_stat(slide, spec: StatSlide, palette: Dict[str, str],
                  resolver: StyleResolver) -> None:
    bg = resolver.bg_color(palette["bg"])
    _set_slide_bg(slide, bg)
    _add_accent_dot(slide, resolver.accent_dot(palette["primary"]),
                    palette["primary"])

    title_style = resolver.style_for("title_style", "title")
    text_color = _on_bg_color(bg, palette)
    _add_styled_textbox(
        slide, Inches(1.0), Inches(0.55), Inches(11.5), Inches(0.8),
        _coerce_styled_text(spec.title), title_style, text_color,
    )

    stat_style = resolver.style_for("stat_style", "stat")
    _add_styled_textbox(
        slide, Inches(0.6), Inches(2.4), Inches(12.1), Inches(2.5),
        _coerce_styled_text(spec.stat), stat_style, palette["primary"],
    )

    caption_style = resolver.style_for("caption_style", "caption")
    _add_styled_textbox(
        slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.8),
        _coerce_styled_text(spec.caption), caption_style, palette["muted"],
    )


RENDERERS = {
    "title": (TitleSlide, _render_title),
    "section": (SectionSlide, _render_section),
    "content": (ContentSlide, _render_content),
    "two_column": (TwoColumnSlide, _render_two_column),
    "stat": (StatSlide, _render_stat),
}


# ---------- Входная схема и сам tool ----------

class PresentationInput(BaseModel):
    title: str = Field(..., description="Название колоды (для метаданных файла)")
    palette: Literal[
        "midnight_executive", "forest_moss", "coral_energy", "warm_terracotta",
        "ocean_gradient", "charcoal_minimal", "teal_trust", "berry_cream", "cherry_bold",
    ] = Field("midnight_executive", description="Готовая палитра. Переопределяется style.custom_palette.")
    style: Optional[PresentationStyle] = Field(
        None,
        description="Глобальные настройки: размер слайда, кастомная палитра, дефолтные стили.",
    )
    slides: List[dict] = Field(
        ...,
        description=(
            "Список слайдов. Каждый — словарь с 'type' (title|section|content|two_column|stat). "
            "Текстовые поля принимают либо строку, либо {text, style}. "
            "Списки буллетов — массив строк или объектов {text, style}. "
            "У каждого слайда может быть поле 'style' (SlideStyle)."
        ),
    )
    output_path: str = Field("presentation.pptx", description="Куда сохранить .pptx")


@tool("create_presentation", args_schema=PresentationInput)
def create_presentation(
    title: str,
    slides: List[dict],
    palette: str = "midnight_executive",
    style: Optional[dict] = None,
    output_path: str = "presentation.pptx",
) -> str:
    """Сгенерировать .pptx-презентацию со свободной кастомизацией стилей.

    Каждое текстовое поле принимает либо строку, либо объект {text, style},
    где style — TextStyle (font, size, bold, italic, color, align, ...).
    Слайды могут иметь поле 'style' (SlideStyle) для переопределения фона
    и стилей блоков. Параметр верхнего уровня 'style' (PresentationStyle)
    задаёт глобальные дефолты и кастомную палитру.

    Если ничего не задано — работают разумные дефолты.
    Возвращает абсолютный путь к созданному файлу.
    """
    # Резолвим палитру: custom_palette > именованная
    pres_style = PresentationStyle(**style) if style else PresentationStyle()
    if pres_style.custom_palette:
        required = {"primary", "secondary", "accent", "bg", "text", "muted"}
        missing = required - set(pres_style.custom_palette.keys())
        if missing:
            raise ValueError(f"custom_palette не хватает ключей: {sorted(missing)}")
        pal = dict(pres_style.custom_palette)
    else:
        if palette not in PALETTES:
            raise ValueError(
                f"Неизвестная палитра '{palette}'. Доступны: {sorted(PALETTES)}"
            )
        pal = dict(PALETTES[palette])

    prs = Presentation()
    prs.slide_width = Inches(pres_style.slide_width_inches)
    prs.slide_height = Inches(pres_style.slide_height_inches)
    blank_layout = prs.slide_layouts[6]

    for idx, raw in enumerate(slides):
        if "type" not in raw:
            raise ValueError(f"Слайд #{idx}: отсутствует поле 'type'")
        stype = raw["type"]
        if stype not in RENDERERS:
            raise ValueError(
                f"Слайд #{idx}: неизвестный type='{stype}'. "
                f"Доступны: {sorted(RENDERERS)}"
            )
        model_cls, render_fn = RENDERERS[stype]
        spec = model_cls(**raw)

        resolver = StyleResolver(
            palette=pal,
            presentation_defaults=pres_style.defaults,
            slide_style=spec.style,
        )

        slide = prs.slides.add_slide(blank_layout)
        render_fn(slide, spec, pal, resolver)

    prs.core_properties.title = title
    out = Path(output_path).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return str(out)


# ---------- Демо ----------

if __name__ == "__main__":
    # 1) Минимальный пример — никакой кастомизации.
    create_presentation.invoke({
        "title": "Базовый пример",
        "palette": "midnight_executive",
        "slides": [
            {"type": "title", "title": "Без кастомизации",
             "subtitle": "Работает на дефолтах"},
            {"type": "content", "title": "Обычные буллеты",
             "bullets": ["Один", "Два", "Три"]},
        ],
        "output_path": "/mnt/user-data/outputs/demo_default.pptx",
    })

    # 2) Полная кастомизация — все три уровня.
    create_presentation.invoke({
        "title": "Кастомный стиль",
        "style": {
            "custom_palette": {
                "primary": "FF6B35", "secondary": "F7C59F", "accent": "FFFFFF",
                "bg": "FFFCF2", "text": "1A1A1A", "muted": "888888",
            },
            "defaults": {
                "title_style": {"font": "Helvetica", "size": 42, "bold": True},
                "body_style": {"font": "Helvetica", "size": 18, "bullet_char": "→"},
            },
        },
        "slides": [
            {
                "type": "title",
                "title": "Свой шрифт, свои цвета",
                "subtitle": {"text": "и даже свой маркер списка",
                              "style": {"italic": True, "size": 24}},
            },
            {
                "type": "content",
                "title": {"text": "Смешанный стиль", "style": {"color": "FF6B35"}},
                "bullets": [
                    "Обычный пункт",
                    {"text": "Жирный важный пункт",
                     "style": {"bold": True, "color": "990011"}},
                    {"text": "Курсивный пункт",
                     "style": {"italic": True, "bullet_char": "★"}},
                    "Снова обычный",
                ],
                "style": {
                    "bg_color": "1A1A1A",
                    "accent_dot": {"enabled": True, "color": "FF6B35",
                                    "diameter_inches": 0.4},
                    "title_style": {"color": "FFFFFF"},
                    "body_style": {"color": "F7C59F"},
                },
            },
            {
                "type": "stat",
                "title": "Огромная цифра",
                "stat": {"text": "99.9%", "style": {"size": 200, "color": "FF6B35"}},
                "caption": "размер задан вручную",
            },
        ],
        "output_path": "/mnt/user-data/outputs/demo_custom.pptx",
    })

    print("Готово: demo_default.pptx и demo_custom.pptx")
